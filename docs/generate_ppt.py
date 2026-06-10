"""
Oil Risk SaaS 발표용 PPT 생성 스크립트
출력: docs/oil_risk_presentation.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.util as util

# ── 색상 팔레트 ──────────────────────────────────────────────
C_BG_DARK   = RGBColor(0x0D, 0x0D, 0x1A)   # 슬라이드 배경 (다크)
C_BG_LIGHT  = RGBColor(0xF5, 0xF7, 0xFA)   # 슬라이드 배경 (라이트)
C_ACCENT    = RGBColor(0x00, 0xA8, 0xFF)   # 포인트 블루
C_GOLD      = RGBColor(0xF5, 0xA6, 0x23)   # 골드
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK_TXT  = RGBColor(0x1A, 0x1A, 0x2E)
C_GRAY      = RGBColor(0x55, 0x55, 0x66)
C_GREEN     = RGBColor(0x27, 0xAE, 0x60)
C_RED       = RGBColor(0xE7, 0x4C, 0x3C)
C_PURPLE    = RGBColor(0x7B, 0x3F, 0xA0)
C_ORANGE    = RGBColor(0xC0, 0x70, 0x20)

SW = Inches(13.33)   # 슬라이드 너비 (와이드)
SH = Inches(7.5)     # 슬라이드 높이


def new_prs():
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]  # 완전 빈 레이아웃
    return prs.slides.add_slide(layout)


def fill_bg(slide, color: RGBColor):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE_TYPE.AUTO_SHAPE if False else 1,  # MSO_AUTO_SHAPE_TYPE.RECTANGLE=1
        left, top, width, height
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, text, left, top, width, height,
                 font_size=Pt(18), bold=False, color=C_WHITE,
                 align=PP_ALIGN.LEFT, wrap=True, font_name='맑은 고딕'):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return txb


def add_para(tf, text, font_size=Pt(16), bold=False, color=C_WHITE,
             align=PP_ALIGN.LEFT, space_before=Pt(6), font_name='맑은 고딕'):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return p


def add_label_box(slide, text, left, top, width, height,
                  bg=C_ACCENT, fg=C_WHITE, font_size=Pt(15), bold=True):
    """배경색 있는 레이블 박스"""
    r = add_rect(slide, left, top, width, height, fill_color=bg)
    tf = r.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = fg
    run.font.name = '맑은 고딕'
    from pptx.util import Pt as pt2
    tf.margin_top    = pt2(8)
    tf.margin_bottom = pt2(8)
    tf.margin_left   = pt2(10)
    tf.margin_right  = pt2(10)
    return r


# ── 슬라이드 생성 함수들 ──────────────────────────────────────

def slide01(prs):
    """훅 — 30초"""
    sl = blank_slide(prs)
    fill_bg(sl, C_BG_DARK)

    # 타임 배지
    add_label_box(sl, '⏱ 30초', Inches(0.3), Inches(0.2), Inches(1.2), Inches(0.4),
                  bg=C_GRAY, font_size=Pt(11))

    # 날짜 대형
    add_text_box(sl, '2026년 4월 7일',
                 Inches(1.0), Inches(1.2), Inches(11.0), Inches(1.1),
                 font_size=Pt(44), bold=True, color=C_GRAY, align=PP_ALIGN.CENTER)

    # 메인 헤드라인
    add_text_box(sl, '유가가 하루 만에 $15 폭락했습니다.',
                 Inches(0.5), Inches(2.2), Inches(12.3), Inches(1.8),
                 font_size=Pt(52), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # 서브 수치
    add_text_box(sl, 'WTI 원유   $71.9  →  $57.0   │   단 하루   │   2026.04.07',
                 Inches(1.0), Inches(4.2), Inches(11.0), Inches(0.7),
                 font_size=Pt(22), bold=False, color=C_RED, align=PP_ALIGN.CENTER)

    # 클로징 멘트
    add_text_box(sl, '"이날 이후, 유가 리스크를 관리하는 기업과 당하는 기업의 운명이 갈렸습니다."',
                 Inches(1.0), Inches(5.4), Inches(11.0), Inches(0.9),
                 font_size=Pt(20), bold=False, color=C_GOLD, align=PP_ALIGN.CENTER)

    # 이미지 플레이스홀더
    r = add_rect(sl, Inches(4.5), Inches(5.0), Inches(4.3), Inches(1.8),
                 fill_color=RGBColor(0x22, 0x22, 0x33),
                 line_color=C_GRAY, line_width=Pt(1))
    add_text_box(sl, '[ 뉴스 헤드라인 스크린샷 삽입 ]',
                 Inches(4.5), Inches(5.5), Inches(4.3), Inches(0.7),
                 font_size=Pt(13), color=C_GRAY, align=PP_ALIGN.CENTER)


def slide02(prs):
    """피해 규모 — 1분"""
    sl = blank_slide(prs)
    fill_bg(sl, C_BG_DARK)

    add_label_box(sl, '⏱ 1분', Inches(0.3), Inches(0.2), Inches(1.0), Inches(0.4),
                  bg=C_GRAY, font_size=Pt(11))

    add_text_box(sl, '유가 $1 차이, 기업엔 수백억입니다',
                 Inches(0.5), Inches(0.7), Inches(12.0), Inches(1.0),
                 font_size=Pt(38), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    cards = [
        ('✈', '대한항공', '연료비 연간 약 3.2조 원', '유가 $1 변동 = 연간 손익 약 300억 차이\n(2024 사업보고서 기반 추정)'),
        ('🚚', '물류·운송', '유류비 = 전체 원가의 30~40%', '유가 10% 상승 → 영업이익 직격'),
        ('🏭', '정유·석화', '원재료 = 원유', '하루 재고 손익이 수십억 단위'),
    ]
    col_w = Inches(3.8)
    for i, (icon, title, sub, body) in enumerate(cards):
        lx = Inches(0.4) + i * (col_w + Inches(0.3))
        ty = Inches(1.9)
        r = add_rect(sl, lx, ty, col_w, Inches(4.5),
                     fill_color=RGBColor(0x1A, 0x1A, 0x35),
                     line_color=C_ACCENT, line_width=Pt(1.5))

        add_text_box(sl, icon, lx, ty + Inches(0.2), col_w, Inches(0.7),
                     font_size=Pt(32), align=PP_ALIGN.CENTER, color=C_WHITE)
        add_text_box(sl, title, lx, ty + Inches(0.95), col_w, Inches(0.6),
                     font_size=Pt(22), bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
        add_text_box(sl, sub, lx, ty + Inches(1.55), col_w, Inches(0.6),
                     font_size=Pt(15), color=C_GOLD, align=PP_ALIGN.CENTER)
        add_text_box(sl, body, lx + Inches(0.1), ty + Inches(2.2), col_w - Inches(0.2), Inches(1.8),
                     font_size=Pt(14), color=RGBColor(0xCC, 0xCC, 0xCC), align=PP_ALIGN.CENTER)

    # 강조 박스
    r2 = add_rect(sl, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.65),
                  fill_color=RGBColor(0x1A, 0x2A, 0x1A),
                  line_color=C_GREEN, line_width=Pt(1))
    add_text_box(sl, '"지금 여러분 회사는 내일 유가가 어디로 갈지 알고 계십니까?"',
                 Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.65),
                 font_size=Pt(17), color=C_GREEN, align=PP_ALIGN.CENTER)


def slide03(prs):
    """기존 방식의 한계 — 1분"""
    sl = blank_slide(prs)
    fill_bg(sl, C_BG_DARK)

    add_label_box(sl, '⏱ 1분', Inches(0.3), Inches(0.2), Inches(1.0), Inches(0.4),
                  bg=C_GRAY, font_size=Pt(11))

    add_text_box(sl, '단일 지표로는 전체 그림이 안 보입니다',
                 Inches(0.5), Inches(0.7), Inches(12.0), Inches(1.0),
                 font_size=Pt(38), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    methods = [
        ('📰', '뉴스만 모니터링',
         '몇 개 기사 읽고\n주관적으로 해석',
         '수급·포지션 데이터와\n통합 불가'),
        ('📊', '전문가 리포트',
         '주간·월간 발행 주기\n매일 바뀌는 시장에 부적합',
         '의사결정 시점과\n항상 어긋남'),
        ('📑', '엑셀 수작업',
         '지표 하나씩 따로 추적\n복합 신호 통합 불가',
         '500개 신호 동시\n처리 불가능'),
    ]
    col_w = Inches(3.8)
    for i, (icon, title, prob, tag) in enumerate(methods):
        lx = Inches(0.4) + i * (col_w + Inches(0.3))
        ty = Inches(1.9)
        r = add_rect(sl, lx, ty, col_w, Inches(4.2),
                     fill_color=RGBColor(0x2A, 0x12, 0x12),
                     line_color=C_RED, line_width=Pt(1.5))

        add_text_box(sl, icon, lx, ty + Inches(0.2), col_w, Inches(0.7),
                     font_size=Pt(34), align=PP_ALIGN.CENTER, color=C_WHITE)
        add_text_box(sl, title, lx, ty + Inches(0.95), col_w, Inches(0.6),
                     font_size=Pt(21), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text_box(sl, prob, lx + Inches(0.1), ty + Inches(1.6), col_w - Inches(0.2), Inches(1.2),
                     font_size=Pt(16), color=RGBColor(0xEE, 0xBB, 0xBB), align=PP_ALIGN.CENTER)
        add_text_box(sl, tag, lx + Inches(0.1), ty + Inches(2.85), col_w - Inches(0.2), Inches(0.8),
                     font_size=Pt(14), bold=True, color=C_RED, align=PP_ALIGN.CENTER)

    add_text_box(sl,
                 '뉴스·리포트·수치 — 각각 따로 보면 항상 그림의 일부만 보입니다',
                 Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.7),
                 font_size=Pt(18), color=C_GOLD, align=PP_ALIGN.CENTER)


def slide04(prs):
    """제품 한 줄 소개 — 30초"""
    sl = blank_slide(prs)
    fill_bg(sl, C_BG_DARK)

    add_label_box(sl, '⏱ 30초', Inches(0.3), Inches(0.2), Inches(1.2), Inches(0.4),
                  bg=C_GRAY, font_size=Pt(11))

    add_text_box(sl, '7일 앞을 먼저 보여드립니다',
                 Inches(0.5), Inches(0.8), Inches(12.0), Inches(1.0),
                 font_size=Pt(42), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # 제품명 뱃지
    r = add_rect(sl, Inches(3.5), Inches(1.9), Inches(6.3), Inches(0.9),
                 fill_color=C_ACCENT)
    add_text_box(sl, '유가 리스크 예측 SaaS',
                 Inches(3.5), Inches(1.95), Inches(6.3), Inches(0.8),
                 font_size=Pt(28), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # 대시보드 플레이스홀더
    r2 = add_rect(sl, Inches(1.5), Inches(3.0), Inches(10.3), Inches(3.6),
                  fill_color=RGBColor(0x10, 0x18, 0x30),
                  line_color=C_ACCENT, line_width=Pt(2))
    add_text_box(sl, '[ 대시보드 스크린샷 삽입 ]\n예측 차트 + 리스크 등급 화면',
                 Inches(1.5), Inches(4.3), Inches(10.3), Inches(1.0),
                 font_size=Pt(16), color=C_GRAY, align=PP_ALIGN.CENTER)

    add_text_box(sl, '브라우저만 있으면 됩니다.  설치 없음.  코딩 없음.',
                 Inches(1.0), Inches(6.8), Inches(11.3), Inches(0.5),
                 font_size=Pt(18), color=C_GOLD, align=PP_ALIGN.CENTER)


def slide05(prs):
    """작동 원리 — 1분 30초"""
    sl = blank_slide(prs)
    fill_bg(sl, C_BG_DARK)

    add_label_box(sl, '⏱ 1분 30초', Inches(0.3), Inches(0.2), Inches(1.6), Inches(0.4),
                  bg=C_GRAY, font_size=Pt(11))

    add_text_box(sl, '기상청처럼 — 500가지 정보를 읽어 7일 후 유가를 예보합니다',
                 Inches(0.5), Inches(0.7), Inches(12.0), Inches(1.0),
                 font_size=Pt(32), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # 입력 박스
    inputs = ['📦 EIA 원유 재고', '📰 뉴스·지정학 감성 분석', '📊 시장 변동성 (OVX)', '🗓 OPEC 생산 · 회의 일정', '🗺 지정학 이슈']
    bx = Inches(0.4)
    by = Inches(1.9)
    bw = Inches(2.6)
    bh = Inches(0.72)
    add_label_box(sl, '[ 입력 ]', bx, by, bw, Inches(0.45),
                  bg=RGBColor(0x1A, 0x2A, 0x3A), fg=C_ACCENT, font_size=Pt(13))
    for j, inp in enumerate(inputs):
        r = add_rect(sl, bx, by + Inches(0.5 + j * 0.8), bw, bh,
                     fill_color=RGBColor(0x18, 0x22, 0x33),
                     line_color=C_ACCENT, line_width=Pt(0.8))
        add_text_box(sl, inp, bx + Inches(0.05), by + Inches(0.5 + j * 0.8), bw - Inches(0.1), bh,
                     font_size=Pt(13), color=C_WHITE, align=PP_ALIGN.CENTER)

    # 화살표 3개
    for j in range(3):
        ar = slide05_arrow(sl, Inches(3.15), Inches(2.55 + j * 0.85))

    # AI 박스
    cx = Inches(3.3)
    cy = Inches(1.9)
    cw = Inches(4.0)
    ch = Inches(5.0)
    r = add_rect(sl, cx, cy, cw, ch,
                 fill_color=RGBColor(0x10, 0x18, 0x35),
                 line_color=C_GOLD, line_width=Pt(2.5))
    add_text_box(sl, '[ AI 분석 ]', cx, cy + Inches(0.1), cw, Inches(0.4),
                 font_size=Pt(13), bold=True, color=C_GOLD, align=PP_ALIGN.CENTER)
    layers = [
        ('AI 예측 모델 (머신러닝)', C_ACCENT),
        ('+', C_GRAY),
        ('추세·계절성 분석', RGBColor(0x27, 0xAE, 0x60)),
        ('+', C_GRAY),
        ('실시간 감성 분석', C_PURPLE),
    ]
    for j, (lbl, lc) in enumerate(layers):
        add_text_box(sl, lbl, cx + Inches(0.1), cy + Inches(0.7 + j * 0.8), cw - Inches(0.2), Inches(0.7),
                     font_size=Pt(16 if lbl != '+' else 20), bold=(lbl != '+'),
                     color=lc, align=PP_ALIGN.CENTER)
    add_text_box(sl, '3단계 복합 분석', cx, cy + Inches(4.4), cw, Inches(0.45),
                 font_size=Pt(13), bold=True, color=C_GOLD, align=PP_ALIGN.CENTER)

    # 화살표
    slide05_arrow(sl, Inches(7.45), Inches(3.5))

    # 출력 박스
    ox = Inches(7.65)
    oy = Inches(1.9)
    ow = Inches(3.8)
    oh = Inches(5.0)
    r2 = add_rect(sl, ox, oy, ow, oh,
                  fill_color=RGBColor(0x0D, 0x22, 0x0D),
                  line_color=C_GREEN, line_width=Pt(2))
    add_text_box(sl, '[ 출력 ]', ox, oy + Inches(0.1), ow, Inches(0.4),
                 font_size=Pt(13), bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)
    outputs = [
        ('📈 7일 예측 가격', C_ACCENT),
        ('⚠️  리스크 등급', C_GOLD),
        ('🔔 극단 변동성 경보', C_RED),
    ]
    for j, (lbl, lc) in enumerate(outputs):
        r3 = add_rect(sl, ox + Inches(0.15), oy + Inches(0.7 + j * 1.3),
                      ow - Inches(0.3), Inches(1.1),
                      fill_color=RGBColor(0x18, 0x30, 0x18),
                      line_color=lc, line_width=Pt(1.2))
        add_text_box(sl, lbl, ox + Inches(0.15), oy + Inches(0.7 + j * 1.3),
                     ow - Inches(0.3), Inches(1.1),
                     font_size=Pt(16), bold=True, color=lc, align=PP_ALIGN.CENTER)

    add_text_box(sl, '매일 새벽 자동 실행  ·  데이터 수집 → 분석 → 알림까지 무인 자동화',
                 Inches(0.5), Inches(7.1), Inches(12.3), Inches(0.4),
                 font_size=Pt(14), color=C_GRAY, align=PP_ALIGN.CENTER)


def slide05_arrow(slide, lx, ly):
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    shp = slide.shapes.add_shape(
        13,  # RIGHT_ARROW
        lx, ly, Inches(0.5), Inches(0.4)
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = C_ACCENT
    shp.line.fill.background()
    return shp


def slide06(prs):
    """핵심 기능 3가지 — 1분"""
    sl = blank_slide(prs)
    fill_bg(sl, C_BG_DARK)

    add_label_box(sl, '⏱ 1분', Inches(0.3), Inches(0.2), Inches(1.0), Inches(0.4),
                  bg=C_GRAY, font_size=Pt(11))

    add_text_box(sl, '담당자가 실제로 쓰는 기능 3가지',
                 Inches(0.5), Inches(0.7), Inches(12.0), Inches(1.0),
                 font_size=Pt(38), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    features = [
        ('📈', '7일 가격 예측', C_ACCENT,
         'D+1 ~ D+7\n일별 예상 가격 제공',
         '"헤지 계약 타이밍\n 결정에 활용"'),
        ('⚠️', '리스크 등급 알림', C_GOLD,
         'SURGE_RISK / DROP_RISK\n/ CAUTION / NORMAL 자동 분류',
         '"오늘 구매할까\n 말까 판단"'),
        ('🚨', '극단 변동성 경보', C_RED,
         '변동성 지수 · 예상 손실\n시장 급등·급락 징후 자동 감지',
         '"최악의 날을\n 대비할 수 있음"'),
    ]
    cw = Inches(3.9)
    for i, (icon, title, color, body, quote) in enumerate(features):
        lx = Inches(0.35) + i * (cw + Inches(0.28))
        ty = Inches(1.85)
        r = add_rect(sl, lx, ty, cw, Inches(5.1),
                     fill_color=RGBColor(0x10, 0x18, 0x30),
                     line_color=color, line_width=Pt(2.5))
        add_text_box(sl, icon, lx, ty + Inches(0.15), cw, Inches(0.8),
                     font_size=Pt(38), align=PP_ALIGN.CENTER, color=C_WHITE)
        add_text_box(sl, title, lx, ty + Inches(1.0), cw, Inches(0.65),
                     font_size=Pt(20), bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text_box(sl, body, lx + Inches(0.1), ty + Inches(1.7), cw - Inches(0.2), Inches(1.2),
                     font_size=Pt(15), color=RGBColor(0xCC, 0xCC, 0xCC), align=PP_ALIGN.CENTER)
        r2 = add_rect(sl, lx + Inches(0.15), ty + Inches(3.1), cw - Inches(0.3), Inches(1.6),
                      fill_color=RGBColor(0x1A, 0x1A, 0x2E),
                      line_color=RGBColor(0x44, 0x44, 0x55), line_width=Pt(1))
        add_text_box(sl, quote, lx + Inches(0.15), ty + Inches(3.1), cw - Inches(0.3), Inches(1.6),
                     font_size=Pt(14), color=C_GOLD, align=PP_ALIGN.CENTER)


def slide06a(prs):
    """데이터 출처 — 45초"""
    sl = blank_slide(prs)
    fill_bg(sl, C_BG_DARK)

    add_label_box(sl, '⏱ 45초', Inches(0.3), Inches(0.2), Inches(1.2), Inches(0.4),
                  bg=C_GRAY, font_size=Pt(11))

    add_text_box(sl, '데이터는 전부 공식 기관 출처입니다',
                 Inches(0.5), Inches(0.7), Inches(12.0), Inches(1.0),
                 font_size=Pt(38), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    rows = [
        ('📦 원유 재고·수급 통계',   'EIA\n(미국 에너지청)',           '미국 연방정부 공식 기관',  C_ACCENT),
        ('💹 시장 포지션 (투기/헤지)', 'CFTC\n(미 상품선물거래위원회)', '미국 금융 규제기관',       RGBColor(0x27, 0xAE, 0x60)),
        ('📰 지정학 뉴스 500+ 소스',  'NewsAPI · RSS',                 '글로벌 주요 언론 실시간',  C_GOLD),
        ('📊 시장 가격·변동성',       'Yahoo Finance',                 '글로벌 시장 표준 데이터',  C_PURPLE),
    ]
    header_y = Inches(1.85)
    headers = ['데이터', '제공 기관', '신뢰 근거']
    col_x = [Inches(0.4), Inches(4.5), Inches(8.3)]
    col_w2 = [Inches(4.0), Inches(3.7), Inches(4.2)]
    for k, (hdr, hx, hw) in enumerate(zip(headers, col_x, col_w2)):
        r = add_rect(sl, hx, header_y, hw, Inches(0.55),
                     fill_color=RGBColor(0x1A, 0x2A, 0x3A))
        add_text_box(sl, hdr, hx, header_y, hw, Inches(0.55),
                     font_size=Pt(16), bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

    row_h = Inches(1.0)
    for i, (data, org, trust, color) in enumerate(rows):
        ry = header_y + Inches(0.6) + i * row_h
        bg_c = RGBColor(0x12, 0x18, 0x28) if i % 2 == 0 else RGBColor(0x18, 0x20, 0x35)
        for k, (val, cx, cw3) in enumerate(zip([data, org, trust], col_x, col_w2)):
            r = add_rect(sl, cx, ry, cw3, row_h - Inches(0.05),
                         fill_color=bg_c,
                         line_color=color if k == 0 else RGBColor(0x33, 0x33, 0x44),
                         line_width=Pt(2) if k == 0 else Pt(0.5))
            add_text_box(sl, val, cx + Inches(0.05), ry, cw3 - Inches(0.1), row_h - Inches(0.05),
                         font_size=Pt(14), color=color if k == 0 else C_WHITE,
                         bold=(k == 0), align=PP_ALIGN.CENTER)

    r3 = add_rect(sl, Inches(0.4), Inches(6.55), Inches(12.5), Inches(0.7),
                  fill_color=RGBColor(0x0D, 0x20, 0x0D),
                  line_color=C_GREEN, line_width=Pt(1.5))
    add_text_box(sl,
                 '골드만삭스 원자재 리서치팀이 참조하는 것과 동일한 공식 데이터를 사용합니다',
                 Inches(0.4), Inches(6.55), Inches(12.5), Inches(0.7),
                 font_size=Pt(16), bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)


def slide06b(prs):
    """한계 공개 — 45초"""
    sl = blank_slide(prs)
    fill_bg(sl, C_BG_DARK)

    add_label_box(sl, '⏱ 45초', Inches(0.3), Inches(0.2), Inches(1.2), Inches(0.4),
                  bg=C_GRAY, font_size=Pt(11))

    add_text_box(sl, '100% 예측은 거짓말입니다',
                 Inches(0.5), Inches(0.65), Inches(12.0), Inches(0.75),
                 font_size=Pt(38), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text_box(sl, '우리는  "언제 조심해야 하는지"를 알려드립니다',
                 Inches(0.5), Inches(1.35), Inches(12.0), Inches(0.6),
                 font_size=Pt(22), bold=False, color=C_GOLD, align=PP_ALIGN.CENTER)

    col_items = [
        ('✅ 잘 작동하는 상황', C_GREEN,
         ['평시 가격 흐름 예측', '방향성 적중률 61%\n(학습 데이터 기준)', '오차 평균 $2 이내\n(평상시 기준)', '가격 방향 전환 사전 감지']),
        ('❌ 솔직한 한계', C_RED,
         ['주말 외부 충격\n(관세 발표·전쟁 등)', '당일 급락\n(전일까지 신호 없음)', '지정학 충격 직후\n72시간', '']),
    ]
    cw = Inches(5.8)
    for i, (title, color, items) in enumerate(col_items):
        lx = Inches(0.5) + i * (cw + Inches(0.8))
        ty = Inches(2.1)
        r = add_rect(sl, lx, ty, cw, Inches(4.6),
                     fill_color=RGBColor(0x12, 0x18, 0x28),
                     line_color=color, line_width=Pt(2.5))
        add_text_box(sl, title, lx, ty + Inches(0.1), cw, Inches(0.6),
                     font_size=Pt(20), bold=True, color=color, align=PP_ALIGN.CENTER)
        for j, item in enumerate(items):
            if item:
                add_text_box(sl, item, lx + Inches(0.15), ty + Inches(0.8 + j * 0.9),
                             cw - Inches(0.3), Inches(0.85),
                             font_size=Pt(15), color=C_WHITE, align=PP_ALIGN.LEFT)

    r2 = add_rect(sl, Inches(0.5), Inches(6.95), Inches(12.3), Inches(0.35),
                  fill_color=RGBColor(0x1A, 0x1A, 0x10),
                  line_color=C_GOLD, line_width=Pt(1))
    add_text_box(sl,
                 '예측 신뢰 구간과 적용 범위를 명시합니다 — 이것이 신뢰의 근거입니다',
                 Inches(0.5), Inches(6.95), Inches(12.3), Inches(0.45),
                 font_size=Pt(15), color=C_GOLD, align=PP_ALIGN.CENTER)


def slide07(prs):
    """성능 검증 — 1분"""
    sl = blank_slide(prs)
    fill_bg(sl, C_BG_DARK)

    add_label_box(sl, '⏱ 1분', Inches(0.3), Inches(0.2), Inches(1.0), Inches(0.4),
                  bg=C_GRAY, font_size=Pt(11))

    add_text_box(sl, '숫자로 증명합니다  —  3년치 실제 데이터 시뮬레이션 기준',
                 Inches(0.5), Inches(0.7), Inches(12.0), Inches(0.9),
                 font_size=Pt(32), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # 비교 테이블
    t_rows = [
        ('Persistence\n(아무것도 안 한 것)', '1.000', '기준선', C_GRAY),
        ('기존 통계 예측 모델\n(ETS / VAR)',       '~1.000', '기준선 수준', C_GRAY),
        ('우리 시스템',                        '0.602',  '기준 대비 40% 개선 ↓', C_GOLD),
    ]
    col_x2 = [Inches(0.5), Inches(5.0), Inches(8.5)]
    col_w3 = [Inches(4.4), Inches(3.4), Inches(4.0)]
    hdrs2 = ['비교 대상', '오차 점수 (낮을수록 좋음)', '설명']
    hy = Inches(1.75)
    for k, (h, cx, cw3) in enumerate(zip(hdrs2, col_x2, col_w3)):
        r = add_rect(sl, cx, hy, cw3, Inches(0.55),
                     fill_color=RGBColor(0x1A, 0x2A, 0x3A))
        add_text_box(sl, h, cx, hy, cw3, Inches(0.55),
                     font_size=Pt(16), bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

    for i, (name, score, note, color) in enumerate(t_rows):
        ry = hy + Inches(0.6) + i * Inches(1.05)
        bg_c = RGBColor(0x0D, 0x20, 0x0D) if i == 2 else RGBColor(0x18, 0x1A, 0x28)
        lw2 = Pt(3) if i == 2 else Pt(0.8)
        lc2 = C_GOLD if i == 2 else RGBColor(0x33, 0x33, 0x44)
        for k, (val, cx, cw3) in enumerate(zip([name, score, note], col_x2, col_w3)):
            r = add_rect(sl, cx, ry, cw3, Inches(1.0),
                         fill_color=bg_c, line_color=lc2, line_width=lw2)
            fs = Pt(28) if (k == 1 and i == 2) else Pt(15)
            fc2 = color if k in (1, 2) else C_WHITE
            bold2 = (k == 1 and i == 2) or (k == 2 and i == 2)
            add_text_box(sl, val, cx + Inches(0.05), ry, cw3 - Inches(0.1), Inches(1.0),
                         font_size=fs, bold=bold2, color=fc2, align=PP_ALIGN.CENTER)

    # 3개 지표
    stats = [('40%', '예측 오차 개선', C_GOLD), ('61%', '방향성 적중률\n(학습 기준)', C_ACCENT), ('$2', '평균 예측 오차\n(평상시 기준)', C_GREEN)]
    sw_each = Inches(3.8)
    for i, (num, label, color) in enumerate(stats):
        lx = Inches(0.5) + i * (sw_each + Inches(0.25))
        ty = Inches(5.0)
        r = add_rect(sl, lx, ty, sw_each, Inches(1.85),
                     fill_color=RGBColor(0x10, 0x18, 0x30),
                     line_color=color, line_width=Pt(2))
        add_text_box(sl, num, lx, ty + Inches(0.1), sw_each, Inches(0.9),
                     font_size=Pt(48), bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text_box(sl, label, lx, ty + Inches(1.05), sw_each, Inches(0.55),
                     font_size=Pt(16), color=C_WHITE, align=PP_ALIGN.CENTER)

    add_text_box(sl, '3년치 실제 데이터로 시뮬레이션 검증  ·  점수: 낮을수록 정확 (1.0 = 예측 안 한 것과 동일)',
                 Inches(0.5), Inches(7.1), Inches(12.3), Inches(0.35),
                 font_size=Pt(12), color=C_GRAY, align=PP_ALIGN.CENTER)


def slide08(prs):
    """타겟 고객 포지셔닝 — 1분 30초"""
    sl = blank_slide(prs)
    fill_bg(sl, C_BG_DARK)

    add_label_box(sl, '⏱ 1분 30초', Inches(0.3), Inches(0.2), Inches(1.6), Inches(0.4),
                  bg=C_GRAY, font_size=Pt(11))

    add_text_box(sl, '전문성은 Bloomberg급  —  가격은 그 이하',
                 Inches(0.5), Inches(0.7), Inches(12.0), Inches(0.9),
                 font_size=Pt(36), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # 포지셔닝 맵 박스
    map_l, map_t = Inches(0.5), Inches(1.7)
    map_w, map_h = Inches(6.0), Inches(4.8)
    r = add_rect(sl, map_l, map_t, map_w, map_h,
                 fill_color=RGBColor(0x10, 0x14, 0x28),
                 line_color=C_GRAY, line_width=Pt(1))

    # 축 레이블
    add_text_box(sl, '고가 ↑', map_l + Inches(2.7), map_t + Inches(0.05), Inches(1.5), Inches(0.4),
                 font_size=Pt(13), color=C_GRAY, align=PP_ALIGN.CENTER)
    add_text_box(sl, '저가 ↓', map_l + Inches(2.7), map_t + Inches(4.3), Inches(1.5), Inches(0.4),
                 font_size=Pt(13), color=C_GRAY, align=PP_ALIGN.CENTER)
    add_text_box(sl, '← 복잡함', map_l, map_t + Inches(2.15), Inches(1.3), Inches(0.4),
                 font_size=Pt(12), color=C_GRAY, align=PP_ALIGN.LEFT)
    add_text_box(sl, '단순함 →', map_l + map_w - Inches(1.3), map_t + Inches(2.15), Inches(1.3), Inches(0.4),
                 font_size=Pt(12), color=C_GRAY, align=PP_ALIGN.RIGHT)

    # 플레이어 점
    players = [
        ('Bloomberg\nTerminal', map_l + Inches(1.0), map_t + Inches(0.7), C_RED),
        ('우리 시스템\n★', map_l + Inches(2.3), map_t + Inches(2.1), C_GOLD),
        ('엑셀 / 직관', map_l + Inches(4.2), map_t + Inches(3.4), C_GRAY),
    ]
    for lbl, px, py, pc in players:
        r2 = add_rect(sl, px, py, Inches(1.5), Inches(0.9),
                      fill_color=RGBColor(0x20, 0x20, 0x40),
                      line_color=pc, line_width=Pt(2))
        add_text_box(sl, lbl, px, py, Inches(1.5), Inches(0.9),
                     font_size=Pt(13), bold=True, color=pc, align=PP_ALIGN.CENTER)

    # 타겟 고객
    targets = [
        ('🚚', '중견 물류·운송사', '유류비 = 원가의 30~40%'),
        ('📈', '에너지 트레이더', '포지션 진입 타이밍'),
        ('🏭', '정유·석화 구매팀', '원재료 구매 최적화'),
    ]
    tx = Inches(7.0)
    add_text_box(sl, '주요 타겟 고객', tx, Inches(1.7), Inches(5.9), Inches(0.5),
                 font_size=Pt(18), bold=True, color=C_ACCENT, align=PP_ALIGN.LEFT)
    for i, (icon, name, desc) in enumerate(targets):
        ty2 = Inches(2.3) + i * Inches(1.5)
        r3 = add_rect(sl, tx, ty2, Inches(5.9), Inches(1.35),
                      fill_color=RGBColor(0x12, 0x18, 0x28),
                      line_color=C_ACCENT, line_width=Pt(1.5))
        add_text_box(sl, icon + '  ' + name, tx + Inches(0.1), ty2 + Inches(0.1),
                     Inches(5.7), Inches(0.6),
                     font_size=Pt(18), bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
        add_text_box(sl, desc, tx + Inches(0.15), ty2 + Inches(0.7),
                     Inches(5.6), Inches(0.5),
                     font_size=Pt(14), color=C_GOLD, align=PP_ALIGN.LEFT)

    add_text_box(sl,
                 '"Bloomberg Terminal은 너무 비쌉니다. 엑셀은 너무 허술합니다. 저희는 그 사이입니다."',
                 Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.6),
                 font_size=Pt(16), color=C_GOLD, align=PP_ALIGN.CENTER)


def slide09(prs):
    """수익 모델 — 1분"""
    sl = blank_slide(prs)
    fill_bg(sl, C_BG_DARK)

    add_label_box(sl, '⏱ 1분', Inches(0.3), Inches(0.2), Inches(1.0), Inches(0.4),
                  bg=C_GRAY, font_size=Pt(11))

    add_text_box(sl, 'Freemium → SaaS 전환 수익 모델',
                 Inches(0.5), Inches(0.7), Inches(12.0), Inches(0.9),
                 font_size=Pt(38), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    plans = [
        ('FREE', '₩ 0 / 월', C_GRAY,
         '체험 · 유입',
         ['기본 예측 조회', '리스크 등급 확인', '', '']),
        ('STANDARD', '₩ 490,000 / 월', C_ACCENT,
         '중소 물류사',
         ['7일 예측 전체', '리스크 알림', '이메일 알림', '']),
        ('PRO', '₩ 1,290,000 / 월', C_GOLD,
         '트레이더 · 대기업',
         ['전체 기능 + 모델별 상세', '극단 변동성 경보', '커스텀 임계값 설정', '성과 분석 리포트']),
    ]
    cw = Inches(3.9)
    for i, (tier, price, color, target, items) in enumerate(plans):
        lx = Inches(0.35) + i * (cw + Inches(0.3))
        ty = Inches(1.85)
        r = add_rect(sl, lx, ty, cw, Inches(4.75),
                     fill_color=RGBColor(0x10, 0x18, 0x30),
                     line_color=color, line_width=Pt(2.5 if i == 2 else 1.5))

        add_label_box(sl, tier, lx, ty, cw, Inches(0.65),
                      bg=color, fg=C_WHITE if i != 0 else C_DARK_TXT,
                      font_size=Pt(20))
        add_text_box(sl, price, lx, ty + Inches(0.7), cw, Inches(0.75),
                     font_size=Pt(20), bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text_box(sl, target, lx, ty + Inches(1.45), cw, Inches(0.5),
                     font_size=Pt(14), color=C_GRAY, align=PP_ALIGN.CENTER)

        for j, item in enumerate(items):
            if item:
                add_text_box(sl, '✓  ' + item, lx + Inches(0.2), ty + Inches(2.1 + j * 0.6),
                             cw - Inches(0.3), Inches(0.5),
                             font_size=Pt(14), color=C_WHITE, align=PP_ALIGN.LEFT)

    r2 = add_rect(sl, Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.55),
                  fill_color=RGBColor(0x10, 0x20, 0x10),
                  line_color=C_GREEN, line_width=Pt(2))
    add_text_box(sl, 'Standard 200명 = 월 1억  ·  Pro 100명 = 월 1.3억  =  연 최대 15.5억',
                 Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.55),
                 font_size=Pt(20), bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)


def slide09a(prs):
    """시장 규모 TAM/SAM/SOM — 1분"""
    sl = blank_slide(prs)
    fill_bg(sl, C_BG_DARK)

    add_label_box(sl, '⏱ 1분', Inches(0.3), Inches(0.2), Inches(1.0), Inches(0.4),
                  bg=C_GRAY, font_size=Pt(11))

    add_text_box(sl, '시장은 충분히 큽니다',
                 Inches(0.5), Inches(0.7), Inches(12.0), Inches(0.9),
                 font_size=Pt(42), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # 동심원 스타일 3단 박스 (세로 깔때기)
    tiers = [
        ('TAM', '전 세계 에너지 리스크 관리 솔루션 시장', '$14.2 Billion (2025)', C_ACCENT, Inches(1.0), Inches(10.8)),
        ('SAM', '국내 유류 다소비 기업 (물류·항공·정유)', '약 4,200개사', RGBColor(0x27, 0xAE, 0x60), Inches(2.0), Inches(8.8)),
        ('SOM', '초기 공략 타겟 (중견 물류사·트레이더)', '200개사 × ₩490,000 = 월 1억 / 연 12억', C_GOLD, Inches(3.0), Inches(6.8)),
    ]
    ty = Inches(1.85)
    for i, (tier, desc, num, color, lx, w) in enumerate(tiers):
        r = add_rect(sl, lx, ty + i * Inches(1.5), w, Inches(1.3),
                     fill_color=RGBColor(0x10, 0x18, 0x30),
                     line_color=color, line_width=Pt(2.5))
        add_label_box(sl, tier, lx, ty + i * Inches(1.5), Inches(1.3), Inches(1.3),
                      bg=color, fg=C_WHITE, font_size=Pt(22))
        add_text_box(sl, desc, lx + Inches(1.4), ty + i * Inches(1.5) + Inches(0.1),
                     w - Inches(1.5), Inches(0.55),
                     font_size=Pt(15), color=C_WHITE, align=PP_ALIGN.LEFT)
        add_text_box(sl, num, lx + Inches(1.4), ty + i * Inches(1.5) + Inches(0.65),
                     w - Inches(1.5), Inches(0.55),
                     font_size=Pt(16), bold=True, color=color, align=PP_ALIGN.LEFT)

    add_text_box(sl,
                 '초기 목표: SOM 200개사 달성 (연 12억)  —  틈새 공략으로 점유율 확장',
                 Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.6),
                 font_size=Pt(18), bold=True, color=C_GOLD, align=PP_ALIGN.CENTER)


def slide10(prs):
    """왜 지금, 왜 우리 — 1분"""
    sl = blank_slide(prs)
    fill_bg(sl, C_BG_DARK)

    add_label_box(sl, '⏱ 1분', Inches(0.3), Inches(0.2), Inches(1.0), Inches(0.4),
                  bg=C_GRAY, font_size=Pt(11))

    add_text_box(sl, '4주 만에 MVP에서 SaaS까지  —  Lean 스프린트 4회',
                 Inches(0.5), Inches(0.7), Inches(12.0), Inches(0.9),
                 font_size=Pt(34), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # 타임라인 이미지 플레이스홀더
    r = add_rect(sl, Inches(0.5), Inches(1.7), Inches(12.3), Inches(4.0),
                 fill_color=RGBColor(0x10, 0x14, 0x28),
                 line_color=C_GRAY, line_width=Pt(1.5))

    # 타임라인 이미지 삽입 시도
    timeline_path = os.path.join(os.path.dirname(__file__), 'timeline_biz.png')
    if os.path.exists(timeline_path):
        sl.shapes.add_picture(timeline_path, Inches(0.5), Inches(1.7), Inches(12.3), Inches(4.0))
    else:
        add_text_box(sl, '[ timeline_biz.png 삽입 ]',
                     Inches(0.5), Inches(3.3), Inches(12.3), Inches(0.7),
                     font_size=Pt(16), color=C_GRAY, align=PP_ALIGN.CENTER)

    stats3 = [
        ('★ 예측 오차 40% 개선', '오차 점수 0.602 확정 (역대 최고)', C_GOLD),
        ('★ 모델 신뢰도 확보', '예측 안정성 27% 향상  불필요 신호 9개 제거', C_ACCENT),
        ('★ SaaS 완성', '구독·결제·인증 완전 자동화 (Ph.26)', C_GREEN),
    ]
    sw2 = Inches(4.0)
    for i, (title, sub, color) in enumerate(stats3):
        lx = Inches(0.5) + i * (sw2 + Inches(0.15))
        r2 = add_rect(sl, lx, Inches(5.85), sw2, Inches(1.3),
                      fill_color=RGBColor(0x10, 0x18, 0x30),
                      line_color=color, line_width=Pt(2))
        add_text_box(sl, title, lx + Inches(0.05), Inches(5.9), sw2 - Inches(0.1), Inches(0.5),
                     font_size=Pt(15), bold=True, color=color, align=PP_ALIGN.LEFT)
        add_text_box(sl, sub, lx + Inches(0.05), Inches(6.4), sw2 - Inches(0.1), Inches(0.6),
                     font_size=Pt(12), color=C_GRAY, align=PP_ALIGN.LEFT)


def slide11(prs):
    """시연 예고 — 15초"""
    sl = blank_slide(prs)
    fill_bg(sl, C_BG_DARK)

    add_label_box(sl, '⏱ 15초', Inches(0.3), Inches(0.2), Inches(1.2), Inches(0.4),
                  bg=C_GRAY, font_size=Pt(11))

    add_text_box(sl, '지금 바로 보여드리겠습니다',
                 Inches(0.5), Inches(2.5), Inches(12.0), Inches(1.2),
                 font_size=Pt(52), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    add_text_box(sl, '실제 운영 중인 시스템  ·  라이브 데이터 기반 시연',
                 Inches(1.0), Inches(4.0), Inches(11.0), Inches(0.8),
                 font_size=Pt(24), color=C_ACCENT, align=PP_ALIGN.CENTER)


def slide12(prs):
    """마무리 CTA — 30초"""
    sl = blank_slide(prs)
    fill_bg(sl, C_BG_DARK)

    add_label_box(sl, '⏱ 30초', Inches(0.3), Inches(0.2), Inches(1.2), Inches(0.4),
                  bg=C_GRAY, font_size=Pt(11))

    add_text_box(sl, '날씨 예보 없이 여행 안 가듯',
                 Inches(0.5), Inches(1.3), Inches(12.0), Inches(0.9),
                 font_size=Pt(36), color=C_GRAY, align=PP_ALIGN.CENTER)
    add_text_box(sl, '유가 예보 없이 사업하지 마세요',
                 Inches(0.5), Inches(2.1), Inches(12.0), Inches(1.1),
                 font_size=Pt(44), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # 3줄 요약
    summaries = [
        ('📊', '7일 유가 예측  ·  리스크 등급  ·  극단 변동성 경보', C_ACCENT),
        ('💰', '₩0 무료 체험  →  ₩490,000/월 Standard', C_GREEN),
        ('🚀', '4주 스프린트  —  지금 바로 사용 가능', C_GOLD),
    ]
    for i, (icon, text, color) in enumerate(summaries):
        ty2 = Inches(3.5) + i * Inches(0.85)
        r = add_rect(sl, Inches(1.5), ty2, Inches(10.3), Inches(0.7),
                     fill_color=RGBColor(0x12, 0x18, 0x28),
                     line_color=color, line_width=Pt(1.5))
        add_text_box(sl, icon + '  ' + text, Inches(1.7), ty2,
                     Inches(10.0), Inches(0.7),
                     font_size=Pt(17), color=color, align=PP_ALIGN.LEFT)

    # 연락처
    add_text_box(sl, '문의 · 베타 신청',
                 Inches(2.0), Inches(6.2), Inches(4.0), Inches(0.5),
                 font_size=Pt(18), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text_box(sl, 'namzxc6339@gmail.com',
                 Inches(2.0), Inches(6.7), Inches(4.0), Inches(0.55),
                 font_size=Pt(16), color=C_ACCENT, align=PP_ALIGN.CENTER)

    # QR 플레이스홀더
    r2 = add_rect(sl, Inches(8.5), Inches(5.8), Inches(2.5), Inches(1.55),
                  fill_color=RGBColor(0x22, 0x22, 0x33),
                  line_color=C_GRAY, line_width=Pt(1))
    add_text_box(sl, '[ QR 코드\n베타 가입 링크 ]',
                 Inches(8.5), Inches(6.1), Inches(2.5), Inches(1.0),
                 font_size=Pt(13), color=C_GRAY, align=PP_ALIGN.CENTER)


# ── 메인 ──────────────────────────────────────────────────────

def main():
    prs = new_prs()

    print('슬라이드 생성 중...')
    slide01(prs)  # 훅
    slide02(prs)  # 피해 규모
    slide03(prs)  # 기존 한계
    slide04(prs)  # 제품 소개
    slide05(prs)  # 작동 원리
    slide06(prs)  # 핵심 기능
    slide06a(prs) # 데이터 출처
    slide06b(prs) # 한계 공개
    slide07(prs)  # 성능 검증
    slide08(prs)  # 타겟·포지셔닝
    slide09(prs)  # 수익 모델
    slide09a(prs) # 시장 규모
    slide10(prs)  # 타임라인
    slide11(prs)  # 시연 예고
    slide12(prs)  # CTA

    out = os.path.join(os.path.dirname(__file__), 'oil_risk_presentation.pptx')
    prs.save(out)
    print(f'저장 완료: {out}')
    print(f'총 슬라이드: {len(prs.slides)}장')


if __name__ == '__main__':
    main()
